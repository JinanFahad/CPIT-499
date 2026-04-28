# =====================================================================
# ppt_builder.py — يبني ملف PowerPoint من JSON (نتيجة ai_pitch_engine)
# الفكرة: نستخدم قالب جاهز فيه placeholders زي {{S2_TITLE}} و {{REV_Y1}}
# ثم نستبدل كل placeholder بالقيمة المناسبة من JSON
# هذي الطريقة تحافظ على التصميم الاحترافي للقالب
# =====================================================================

import os
from pptx import Presentation

BASE_DIR = os.path.dirname(os.path.abspath(__file__))


def _replace_in_shape(shape, mapping: dict):
    """يستبدل النصوص داخل run واحد (للحالات البسيطة)"""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for k, v in mapping.items():
                if k in run.text:
                    run.text = run.text.replace(k, str(v))


def _fmt(val):
    """تنسيق الأرقام: 50000 → '50,000 SAR' (للقيم المالية)"""
    try:
        return f"{int(float(str(val).replace(',', '').replace('SAR', '').strip())):,} SAR"
    except (ValueError, TypeError):
        return str(val) if val is not None else ""


def _build_mapping_from_deck(deck: dict) -> dict:
    """يبني قاموس استبدال {placeholder: value} من JSON الـ deck"""
    slides = deck.get("slides", []) or []
    mapping = {}

    def s(i):
        return slides[i] if i < len(slides) else {}

    def trim(text, limit=85):
        if not text or len(text) <= limit:
            return text or ""
        return text[:limit].rsplit(" ", 1)[0] + "…"

    # Slide 1 — Cover
    mapping["{{DECK_TITLE}}"] = deck.get("deck_title", "")
    mapping["{{TAGLINE}}"]    = deck.get("tagline", "")

    # Slides 2–5 — generic
    for i in range(1, 5):
        slide = s(i)
        idx = i + 1
        mapping[f"{{{{S{idx}_TITLE}}}}"]    = slide.get("title", "")
        mapping[f"{{{{S{idx}_SUBTITLE}}}}"] = trim(slide.get("subtitle", ""), 120)
        bullets = slide.get("bullets", [])
        for b in range(3):
            text = bullets[b] if b < len(bullets) else ""
            mapping[f"{{{{S{idx}_B{b+1}}}}}"] = trim(text, 85)

    # Slide 6 — Financial Highlights
    slide6 = s(5)
    mapping["{{S6_TITLE}}"] = slide6.get("title", "")
    nums = slide6.get("numbers", [])
    mapping["{{REV_Y1}}"] = _fmt(nums[0]["value"]) if len(nums) > 0 else ""
    mapping["{{REV_Y2}}"] = _fmt(nums[1]["value"]) if len(nums) > 1 else ""
    mapping["{{REV_Y3}}"] = _fmt(nums[2]["value"]) if len(nums) > 2 else ""

    # Slide 7 — Competitive Advantage
    slide7 = s(6)
    mapping["{{S7_TITLE}}"]            = slide7.get("title", "")
    mapping["{{LIST_OF_COMPETITORS}}"] = "\n".join(slide7.get("bullets", []))

    # Slide 8 — Investment Ask
    slide8 = s(7)
    mapping["{{S9_TITLE}}"] = slide8.get("title", "")

    nums8 = slide8.get("numbers", [])
    # نحسب الإجمالي من مجموع البنود لو funding_needed مو متوفر صريحاً
    total_from_items = 0
    for item in nums8:
        try:
            total_from_items += int(float(str(item.get("value", "0")).replace(",", "").replace("SAR", "").strip()))
        except (ValueError, TypeError):
            pass
    funding = deck.get("funding_needed") or total_from_items
    mapping["{{TOTAL_AMOUNT}}"] = _fmt(funding) if funding else slide8.get("subtitle", "")

    for n in range(3):
        item = nums8[n] if n < len(nums8) else {}
        mapping[f"{{{{B{n+1}_LABEL}}}}"]  = item.get("label", "")
        mapping[f"{{{{B{n+1}_AMOUNT}}}}"] = _fmt(item.get("value", ""))

    return mapping


def _replace_full_paragraph(shape, mapping):
    """
    يدمج كل runs الفقرة الواحدة قبل الاستبدال.
    ضروري لأن PowerPoint أحياناً يقسم {{PLACEHOLDER}} على عدة runs،
    فلو استبدلنا كل run منفصل ما يلقى المتغير كاملاً.
    """
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        full = "".join(r.text for r in paragraph.runs)
        replaced = full
        for k, v in mapping.items():
            replaced = replaced.replace(k, str(v))
        if replaced != full and paragraph.runs:
            for r in paragraph.runs:
                r.text = ""
            paragraph.runs[0].text = replaced


def build_pptx(deck: dict, out_path: str, template_path: str = None):
    """
    يأخذ JSON من ai_pitch_engine ويبني ملف .pptx
    deck: {deck_title, tagline, slides: [...]}
    out_path: مسار الملف النهائي
    template_path: قالب PowerPoint الجاهز (default: templates/pitch_template.pptx)
    """
    if template_path is None:
        template_path = os.path.join(BASE_DIR, "templates", "pitch_template.pptx")

    slides = deck.get("slides", []) or []
    if len(slides) != 8:
        raise ValueError(f"Pitch deck must have exactly 8 slides, got {len(slides)}")

    # نفتح القالب ونبني خريطة الاستبدالات
    prs = Presentation(template_path)
    mapping = _build_mapping_from_deck(deck)

    # نمر على كل شريحة وكل عنصر فيها ونستبدل النصوص
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, mapping)
            _replace_full_paragraph(shape, mapping)

    prs.save(out_path)